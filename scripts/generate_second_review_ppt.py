#!/usr/bin/env python3
"""
Generate Second Review presentation (max 15 slides, 6+ bullets per required topic).

Install: pip install python-pptx

Run from project root:
  python scripts/generate_second_review_ppt.py
  python scripts/generate_second_review_ppt.py -o Second_Review_Generated.pptx
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Content: eight Second Review topics — each has at least six bullet points.
# Edit PROJECT_* and SLIDES below to match your report and screenshots.
# ---------------------------------------------------------------------------

PROJECT_TITLE = "Hospital HRM — Health Record Management"
PROJECT_SUBTITLE = "Blockchain-anchored medical records | Second Review"

SLIDES: list[tuple[str, list[str]]] = [
    (
        "1. Implementation",
        [
            "Stack: Node.js/Express REST API, MongoDB (Mongoose), React (Vite) SPA, JWT auth with role guards.",
            "Medical record integrity: SHA-256 hashing plus optional Ethereum anchoring via self-transaction calldata.",
            "File uploads served under /uploads; JSON body limit tuned for clinical documents and attachments.",
            "CORS and CLIENT_ORIGIN configured for local dev; dotenv-driven secrets (MongoDB, JWT, RPC, keys).",
            "Activity logging and access-control routes support audit trails and delegated patient access.",
            "End-to-end flow: authenticate → CRUD patients/records/vitals/prescriptions → blockchain verify where enabled.",
        ],
    ),
    (
        "2. Modules",
        [
            "Client: Vite + React UI, API service layer, doctor–patient helpers, block explorer utilities.",
            "Server core: Express app, MongoDB connection, centralized error handler, health check endpoint.",
            "Domain routes: auth, users, patients, emergency, records, vitals, prescriptions, reminders, appointments.",
            "Supporting routes: access control, blockchain metadata, uploads, structured activity logs.",
            "Models: User, Patient, MedicalRecord, Vital, Prescription, Appointment, MedicineReminder, AccessControl, ActivityLog, BlockChain, Counter.",
            "Utilities: JWT tokens, PDF/QR generation, patient ID, hash generation, Ethereum anchor, log activity.",
        ],
    ),
    (
        "3. Module Description",
        [
            "Auth module: registration/login, optional auth for public endpoints, role-based guards for clinicians vs admin.",
            "Patient module: demographics, emergency profile, linkage to records, vitals, and appointment history.",
            "Records & vitals: structured clinical data with hashing; prescriptions and reminders for medication adherence.",
            "Blockchain module: stores chain metadata; anchors record hashes on-chain when RPC and wallet are configured.",
            "Access module: grants/revokes delegated access; logs module captures user actions for compliance review.",
            "Uploads & PDF/QR: attachments for records; printable/ scannable artifacts where the UI integrates them.",
        ],
    ),
    (
        "4. Sample Coding / Sample Output / Sample Model",
        [
            "Sample code: ethereumAnchor.js — normalize 64-hex SHA-256, JsonRpcProvider + Wallet, sendTransaction with hash as data.",
            "Sample output: anchor API returns { txHash, chainId } on success; null when env not set or invalid hash.",
            "Sample model: MedicalRecord (and related) fields persisted in MongoDB; indexes and refs as defined in Mongoose schemas.",
            "Sample API: POST/GET under /api/records, /api/blockchain — JSON request/response shapes match client api.js wrappers.",
            "Sample hash flow: hashGenerator produces digest; chain.js / BlockChain model link blocks for tamper-evident history.",
            "Sample client: services/api.js centralizes fetch with credentials for authenticated dashboards and forms.",
        ],
    ),
    (
        "5. Sample Screen Shots",
        [
            "Login / registration and role-appropriate dashboard after successful JWT session.",
            "Patient list or profile view showing demographics and emergency information.",
            "Medical record create/view screen with attachments or hash display where implemented.",
            "Vitals or prescriptions UI demonstrating data entry and saved state confirmation.",
            "Blockchain or verification screen (explorer link / tx hash) after anchoring a record.",
            "Access control or activity log view illustrating audit visibility for reviewers.",
        ],
    ),
    (
        "6. Rough Copy of Full Document",
        [
            "Chapter 1: Introduction — problem statement, objectives, scope of Hospital HRM.",
            "Chapter 2: Literature / related work — EHR standards, security, blockchain in healthcare (summary).",
            "Chapter 3: Requirements — functional (users, patients, records) and non-functional (security, performance).",
            "Chapter 4: System design — architecture diagram, module diagram, database ER-style collection overview.",
            "Chapter 5: Implementation — environment setup, key algorithms (hashing, anchoring), API catalogue.",
            "Chapter 6: Testing, results, conclusion, future work — align narrative with screenshots and demo script.",
        ],
    ),
    (
        "7. Complete Presentation",
        [
            "Slide deck covers all eight review items within the prescribed slide limit (title + eight content slides).",
            "Consistent fonts and hierarchy: title, bullets, optional speaker notes in your rehearsal copy.",
            "Figures numbered to match the report; reference appendix for long code listings.",
            "Transition from architecture → modules → live demo to satisfy examiner flow.",
            "Backup PDF export prepared in case of projector or font substitution issues.",
            "Q&A slide or closing slide with project title, team names, and repository / deployment pointer.",
        ],
    ),
    (
        "8. Full Project Demo",
        [
            "Start stack: MongoDB running, npm run dev (server + client), verify /api/health returns { ok: true }.",
            "Demo path 1: register/login as clinician → create or open patient → add record / vital / prescription.",
            "Demo path 2: show record hash and optional Ethereum anchor (testnet) with transaction hash evidence.",
            "Demo path 3: delegated access grant/revoke and a sample entry in activity logs.",
            "Demo path 4: appointment or reminder workflow relevant to your configured UI pages.",
            "Close with limitations (e.g., testnet only), security notes, and planned enhancements.",
        ],
    ),
]


def _set_slide_title(slide, text: str) -> None:
    if not slide.shapes.title:
        return
    slide.shapes.title.text = text
    tf = slide.shapes.title.text_frame
    for p in tf.paragraphs:
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)


def _fill_body(slide, bullets: list[str]) -> None:
    if len(slide.placeholders) < 2:
        return
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(14)
        p.space_after = Pt(6)
        p.alignment = PP_ALIGN.LEFT


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide (layout 0 is typically title)
    title_layout = prs.slide_layouts[0]
    slide0 = prs.slides.add_slide(title_layout)
    slide0.shapes.title.text = "SECOND REVIEW"
    slide0.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    slide0.shapes.title.text_frame.paragraphs[0].font.bold = True
    if slide0.placeholders[1]:
        sub = slide0.placeholders[1]
        sub.text = f"{PROJECT_TITLE}\n{PROJECT_SUBTITLE}"

    # Content slides — title + bullets (layout 1)
    bullet_layout = prs.slide_layouts[1]
    for title, points in SLIDES:
        s = prs.slides.add_slide(bullet_layout)
        _set_slide_title(s, title)
        _fill_body(s, points)

    return prs


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate Second Review PowerPoint.")
    parser.add_argument(
        "-o",
        "--output",
        type=Path,
        default=Path(__file__).resolve().parent.parent / "Second_Review_Generated.pptx",
        help="Output .pptx path (default: project root Second_Review_Generated.pptx)",
    )
    args = parser.parse_args()

    prs = build_presentation()
    out: Path = args.output
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    n = len(prs.slides)
    print(f"Saved {n} slides to {out} (limit 15: OK).")


if __name__ == "__main__":
    main()
